import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    # Paleta de colores MODO LÍDER
    C_DARK_NAVY = RGBColor(14, 18, 27)      # #0E121B (fondo oscuro de lujo)
    C_SURFACE_DARK = RGBColor(22, 26, 37)   # #161A25
    C_LIGHT_BG = RGBColor(248, 249, 251)    # #F8F9FB (fondo claro limpio y moderno)
    C_CARD_WHITE = RGBColor(255, 255, 255)  # #FFFFFF
    C_GOLD = RGBColor(212, 175, 55)         # #D4AF37
    C_GOLD_DARK = RGBColor(180, 135, 30)    # #B4871E
    C_GOLD_BG = RGBColor(254, 249, 235)     # #FEF9EB
    C_CORAL = RGBColor(233, 69, 96)         # #E94560
    C_CORAL_BG = RGBColor(255, 243, 245)    # #FFF3F5
    C_TEXT_DARK = RGBColor(27, 31, 42)      # #1B1F2A
    C_TEXT_MUTED = RGBColor(90, 100, 115)   # #5A6473
    C_TEXT_LIGHT = RGBColor(245, 245, 247)  # #F5F5F7
    C_BORDER_LIGHT = RGBColor(226, 230, 238)# #E2E6EE

    FONT_MAIN = "Montserrat"

    def set_bg(slide, color):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = color
        bg.line.fill.background()
        return bg

    def add_header(slide, title_text, category_text="CAPÍTULO 1 • EMPRENDER CONSCIENTEMENTE", is_dark=False):
        # Header category
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.35))
        tf_c = cat_box.text_frame
        tf_c.word_wrap = True
        tf_c.margin_left = tf_c.margin_top = tf_c.margin_right = tf_c.margin_bottom = 0
        p_c = tf_c.paragraphs[0]
        p_c.text = category_text.upper()
        p_c.font.name = FONT_MAIN
        p_c.font.size = Pt(10.5)
        p_c.font.bold = True
        p_c.font.color.rgb = C_GOLD if is_dark else C_GOLD_DARK

        # Main title
        t_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(11.733), Inches(0.85))
        tf_t = t_box.text_frame
        tf_t.word_wrap = True
        tf_t.margin_left = tf_t.margin_top = tf_t.margin_right = tf_t.margin_bottom = 0
        p_t = tf_t.paragraphs[0]
        p_t.text = title_text
        p_t.font.name = FONT_MAIN
        p_t.font.size = Pt(21)
        p_t.font.bold = True
        p_t.font.color.rgb = C_TEXT_LIGHT if is_dark else C_TEXT_DARK

    # ==========================================
    # SLIDE 1: PORTADA (OSCURA)
    # ==========================================
    s1 = prs.slides.add_slide(blank_layout)
    set_bg(s1, C_DARK_NAVY)

    # Accent bar
    dec = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.2), Inches(0.12), Inches(3.2))
    dec.fill.solid()
    dec.fill.fore_color.rgb = C_GOLD
    dec.line.fill.background()

    tb1 = s1.shapes.add_textbox(Inches(1.2), Inches(2.1), Inches(11.0), Inches(0.5))
    p = tb1.text_frame.paragraphs[0]
    p.text = "CURSO: CREAR UN NEGOCIO DE CERO A IMPACTO Y VENTAS"
    p.font.name = FONT_MAIN
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = C_GOLD

    tb2 = s1.shapes.add_textbox(Inches(1.2), Inches(2.6), Inches(11.0), Inches(1.8))
    p = tb2.text_frame.paragraphs[0]
    p.text = "CAPITULO 1\nEMPRENDER CONSCIENTEMENTE"
    p.font.name = FONT_MAIN
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = C_TEXT_LIGHT

    tb3 = s1.shapes.add_textbox(Inches(1.2), Inches(4.5), Inches(11.0), Inches(0.8))
    p = tb3.text_frame.paragraphs[0]
    p.text = "MÉTODO MODO LÍDER"
    p.font.name = FONT_MAIN
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(180, 190, 205)

    # ==========================================
    # SLIDE 2: PREGUNTA Y DIFERENCIA CLAVE (CLARA)
    # ==========================================
    s2 = prs.slides.add_slide(blank_layout)
    set_bg(s2, C_LIGHT_BG)
    add_header(s2, "¿Qué diferencia a un emprendedor consciente de otro que no lo es?", "La diferencia clave podría ser:")

    # Tarjeta 1: No consciente
    c1 = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.75), Inches(5.6), Inches(5.0))
    c1.fill.solid()
    c1.fill.fore_color.rgb = C_CARD_WHITE
    c1.line.color.rgb = C_CORAL
    
    tf_c1 = c1.text_frame
    tf_c1.word_wrap = True
    tf_c1.margin_top = Inches(0.3)
    tf_c1.margin_left = tf_c1.margin_right = Inches(0.35)
    
    p = tf_c1.paragraphs[0]
    p.text = "Un emprendedor no consciente pregunta:\n"
    p.font.name = FONT_MAIN
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = C_CORAL

    p2 = tf_c1.add_paragraph()
    p2.text = "“¿Cómo logro que mi negocio funcione?”"
    p2.font.name = FONT_MAIN
    p2.font.size = Pt(15)
    p2.font.bold = True
    p2.font.color.rgb = C_TEXT_DARK

    # Tarjeta 2: Consciente
    c2 = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.75), Inches(5.733), Inches(5.0))
    c2.fill.solid()
    c2.fill.fore_color.rgb = C_CARD_WHITE
    c2.line.color.rgb = C_GOLD
    
    tf_c2 = c2.text_frame
    tf_c2.word_wrap = True
    tf_c2.margin_top = Inches(0.3)
    tf_c2.margin_left = tf_c2.margin_right = Inches(0.35)
    
    p = tf_c2.paragraphs[0]
    p.text = "Un emprendedor consciente pregunta:\n"
    p.font.name = FONT_MAIN
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = C_GOLD_DARK

    p2 = tf_c2.add_paragraph()
    p2.text = "“¿Qué negocio quiero construir y en quién me quiero convertir mientras lo construyo?”"
    p2.font.name = FONT_MAIN
    p2.font.size = Pt(14)
    p2.font.bold = True
    p2.font.color.rgb = C_TEXT_DARK

    p3 = tf_c2.add_paragraph()
    p3.text = "\ny luego se pregunta otras cosas, como con qué estilo de clientes me gustaría trabajar . ¿Cómo logro comunicarme con ellos? ¿Cómo conecto con sus intereses para crear conexión real? Para luego pensar, ¿qué cosas les gustan y les disgustan a mis clientes objetivos? ¿Qué odian, qué temen, que necesitan y de cuáles de sus necesidades voy a ocuparme de resolver?"
    p3.font.name = FONT_MAIN
    p3.font.size = Pt(11)
    p3.font.color.rgb = C_TEXT_DARK

    # ==========================================
    # SLIDE 3: LA DIFERENCIA ESTÁ EN LA MANERA DE PENSAR (CLARA)
    # ==========================================
    s3 = prs.slides.add_slide(blank_layout)
    set_bg(s3, C_LIGHT_BG)
    add_header(s3, "🌱 La diferencia no está en el negocio, sino en la manera de pensar y de emprender tu negocio")

    # Card A: Ambición y Éxito
    box_a = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.733), Inches(1.5))
    box_a.fill.solid()
    box_a.fill.fore_color.rgb = C_CARD_WHITE
    box_a.line.color.rgb = C_BORDER_LIGHT
    tf_a = box_a.text_frame
    tf_a.word_wrap = True
    tf_a.margin_left = tf_a.margin_right = Inches(0.4)
    p = tf_a.paragraphs[0]
    p.text = "Un emprendedor consciente también quiere ganar dinero, crecer, vender y tener éxito. No renuncia a la ambición."
    p.font.name = FONT_MAIN
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = C_GOLD_DARK

    # Card B: Éxito + Impacto + Valores
    box_b = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(3.5), Inches(11.733), Inches(1.5))
    box_b.fill.solid()
    box_b.fill.fore_color.rgb = C_CARD_WHITE
    box_b.line.color.rgb = C_BORDER_LIGHT
    tf_b = box_b.text_frame
    tf_b.word_wrap = True
    tf_b.margin_left = tf_b.margin_right = Inches(0.4)
    p = tf_b.paragraphs[0]
    p.text = "La diferencia es que no separa el éxito económico del impacto que genera, de sus valores y de la persona que está siendo."
    p.font.name = FONT_MAIN
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = C_TEXT_DARK

    # Card C: Trascender miedos
    box_c = s3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.2), Inches(11.733), Inches(1.7))
    box_c.fill.solid()
    box_c.fill.fore_color.rgb = C_CARD_WHITE
    box_c.line.color.rgb = C_BORDER_LIGHT
    tf_c = box_c.text_frame
    tf_c.word_wrap = True
    tf_c.margin_left = tf_c.margin_right = Inches(0.4)
    p = tf_c.paragraphs[0]
    p.text = "Por otro lado un emprendedor consciente sabe que para poder lograr lo que desea necesariamente necesita trascender sus miedos, desactivar sus limitaciones mentales y/o reconocer y resolver las heridas que bloquean su camino y limitaciones hacia su éxito."
    p.font.name = FONT_MAIN
    p.font.size = Pt(12)
    p.font.color.rgb = C_TEXT_DARK

    # ==========================================
    # SLIDE 4: CITA DESTACADA (CLARA CON ORO)
    # ==========================================
    s4 = prs.slides.add_slide(blank_layout)
    set_bg(s4, C_GOLD_BG)

    frame4 = s4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.2), Inches(1.2), Inches(10.933), Inches(5.1))
    frame4.fill.solid()
    frame4.fill.fore_color.rgb = C_CARD_WHITE
    frame4.line.color.rgb = C_GOLD
    frame4.line.width = Pt(2)

    tf_4 = frame4.text_frame
    tf_4.word_wrap = True
    tf_4.margin_top = Inches(1.1)
    tf_4.margin_left = tf_4.margin_right = Inches(0.8)

    p_tag = tf_4.paragraphs[0]
    p_tag.text = "PRINCIPIO FUNDAMENTAL\n"
    p_tag.font.name = FONT_MAIN
    p_tag.font.size = Pt(12)
    p_tag.font.bold = True
    p_tag.alignment = PP_ALIGN.CENTER
    p_tag.font.color.rgb = C_GOLD_DARK

    p_quote = tf_4.add_paragraph()
    p_quote.text = "“Ser un emprendedor consciente no significa emprender sin ambición. Significa tener ambición sin perderte a vos misma en el camino.”"
    p_quote.font.name = FONT_MAIN
    p_quote.font.size = Pt(22)
    p_quote.font.bold = True
    p_quote.alignment = PP_ALIGN.CENTER
    p_quote.font.color.rgb = C_TEXT_DARK

    # ==========================================
    # SLIDE 5: TABLA COMPARATIVA PARTE 1 (CLARA)
    # ==========================================
    s5 = prs.slides.add_slide(blank_layout)
    set_bg(s5, C_LIGHT_BG)
    add_header(s5, "Emprendedor tradicional / no consciente vs. Emprendedor consciente", "CUADRO COMPARATIVO • PARTE 1")

    tbl_s5 = s5.shapes.add_table(6, 2, Inches(0.8), Inches(1.75), Inches(11.733), Inches(5.0))
    t5 = tbl_s5.table
    t5.columns[0].width = Inches(5.866)
    t5.columns[1].width = Inches(5.866)

    # Headers
    h0 = t5.cell(0, 0)
    h0.text = "Emprendedor tradicional / no consciente"
    h0.fill.solid()
    h0.fill.fore_color.rgb = C_CORAL_BG
    for p in h0.text_frame.paragraphs:
        p.font.name = FONT_MAIN
        p.font.size = Pt(11.5)
        p.font.bold = True
        p.font.color.rgb = C_CORAL

    h1 = t5.cell(0, 1)
    h1.text = "Emprendedor consciente"
    h1.fill.solid()
    h1.fill.fore_color.rgb = C_GOLD_BG
    for p in h1.text_frame.paragraphs:
        p.font.name = FONT_MAIN
        p.font.size = Pt(11.5)
        p.font.bold = True
        p.font.color.rgb = C_GOLD_DARK

    rows_5 = [
        ("Busca resultados rápidos", "Construye relaciones y conexión hacia resultados sostenibles"),
        ("Decide principalmente por dinero", "Decide considerando propósito + impacto hacia el dinero"),
        ("Vende lo que puede", "Busca resolver problemas reales de los clientes que decidió atender"),
        ("Compite permanentemente", "Puede colaborar y crear redes desde el respeto, los límites y los acuerdos"),
        ("Trabaja hasta agotarse", "Aprende a gestionar su bienestar, su energía y recursos")
    ]

    for r_idx, (col0, col1) in enumerate(rows_5, start=1):
        c0 = t5.cell(r_idx, 0)
        c0.text = col0
        c0.fill.solid()
        c0.fill.fore_color.rgb = C_CARD_WHITE
        for p in c0.text_frame.paragraphs:
            p.font.name = FONT_MAIN
            p.font.size = Pt(11)
            p.font.color.rgb = C_TEXT_DARK

        c1 = t5.cell(r_idx, 1)
        c1.text = col1
        c1.fill.solid()
        c1.fill.fore_color.rgb = RGBColor(255, 253, 248)
        for p in c1.text_frame.paragraphs:
            p.font.name = FONT_MAIN
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = C_TEXT_DARK

    # ==========================================
    # SLIDE 6: TABLA COMPARATIVA PARTE 2 (CLARA)
    # ==========================================
    s6 = prs.slides.add_slide(blank_layout)
    set_bg(s6, C_LIGHT_BG)
    add_header(s6, "Emprendedor tradicional / no consciente vs. Emprendedor consciente", "CUADRO COMPARATIVO • PARTE 2")

    tbl_s6 = s6.shapes.add_table(6, 2, Inches(0.8), Inches(1.75), Inches(11.733), Inches(5.0))
    t6 = tbl_s6.table
    t6.columns[0].width = Inches(5.866)
    t6.columns[1].width = Inches(5.866)

    h0 = t6.cell(0, 0)
    h0.text = "Emprendedor tradicional / no consciente"
    h0.fill.solid()
    h0.fill.fore_color.rgb = C_CORAL_BG
    for p in h0.text_frame.paragraphs:
        p.font.name = FONT_MAIN
        p.font.size = Pt(11.5)
        p.font.bold = True
        p.font.color.rgb = C_CORAL

    h1 = t6.cell(0, 1)
    h1.text = "Emprendedor consciente"
    h1.fill.solid()
    h1.fill.fore_color.rgb = C_GOLD_BG
    for p in h1.text_frame.paragraphs:
        p.font.name = FONT_MAIN
        p.font.size = Pt(11.5)
        p.font.bold = True
        p.font.color.rgb = C_GOLD_DARK

    rows_6 = [
        ("Busca demostrar que puede", "Busca construir algo que tenga sentido y que genere dinero"),
        ("Se identifica con el resultado", "Aprende del resultado y sigue mejorando o escalando"),
        ("Tiene miedo de equivocarse", "Usa el error como información y crecimiento"),
        ("Entiende que puede crecer a cualquier costo porque así es como se gana dinero", "Define qué costo no está dispuesto a pagar con sabiduría y en coherencia con sus valores"),
        ("Pregunta “¿cuánto gano?”", "Pregunta “¿Qué genero y qué construyo y cuánto gano?”")
    ]

    for r_idx, (col0, col1) in enumerate(rows_6, start=1):
        c0 = t6.cell(r_idx, 0)
        c0.text = col0
        c0.fill.solid()
        c0.fill.fore_color.rgb = C_CARD_WHITE
        for p in c0.text_frame.paragraphs:
            p.font.name = FONT_MAIN
            p.font.size = Pt(11)
            p.font.color.rgb = C_TEXT_DARK

        c1 = t6.cell(r_idx, 1)
        c1.text = col1
        c1.fill.solid()
        c1.fill.fore_color.rgb = RGBColor(255, 253, 248)
        for p in c1.text_frame.paragraphs:
            p.font.name = FONT_MAIN
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = C_TEXT_DARK

    # ==========================================
    # SLIDE 7: PREGUNTAS INCÓMODAS (CLARA)
    # ==========================================
    s7 = prs.slides.add_slide(blank_layout)
    set_bg(s7, C_LIGHT_BG)
    add_header(s7, "🔥 Y después haría una pregunta incómoda", "Les preguntaría:")

    incomodas = [
        ("¿Qué estás dispuesta a hacer para que tu emprendimiento tenga éxito?", "Aquí es donde se crea la mentalidad que necesitás para llegar lejos", C_GOLD_DARK),
        ("“¿Y qué NO estás dispuesta a hacer?”", "Aquí es donde aparecen tus valores, tus límites, tus criterios de decisión básicos", C_CORAL),
        ("¿Cuánto necesitás cobrar por tus productos y/o servicios para que tu propia economía + tu negocio funcione?", "", C_DARK_NAVY)
    ]

    for idx, (pregunta, aclaracion, color_borde) in enumerate(incomodas):
        y_pos = Inches(1.75) + idx * Inches(1.7)
        box = s7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), y_pos, Inches(11.733), Inches(1.5))
        box.fill.solid()
        box.fill.fore_color.rgb = C_CARD_WHITE
        box.line.color.rgb = color_borde
        box.line.width = Pt(1.5)

        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.2)
        tf.margin_left = tf.margin_right = Inches(0.4)

        p = tf.paragraphs[0]
        p.text = pregunta
        p.font.name = FONT_MAIN
        p.font.size = Pt(13.5)
        p.font.bold = True
        p.font.color.rgb = C_TEXT_DARK

        if aclaracion:
            p2 = tf.add_paragraph()
            p2.text = aclaracion
            p2.font.name = FONT_MAIN
            p2.font.size = Pt(11.5)
            p2.font.color.rgb = C_TEXT_MUTED

    # ==========================================
    # SLIDE 8: LA FÓRMULA (OSCURA DE ALTO IMPACTO)
    # ==========================================
    s8 = prs.slides.add_slide(blank_layout)
    set_bg(s8, C_DARK_NAVY)
    add_header(s8, "EMPRENDER CONSCIENTEMENTE = PROPÓSITO + IDENTIDAD + MENTALIDAD + ACCIÓN CON ESTRATEGIA = IMPACTO + DINERO QUE FLUYE HACIA TU NEGOCIO", "LA FÓRMULA", is_dark=True)

    items_formula = [
        ("• Propósito:", "¿para qué?"),
        ("• Identidad:", "¿quién soy y qué quiero construir?"),
        ("• Mentalidad:", "¿qué necesito transformar?"),
        ("• Acción:", "¿qué hago?"),
        ("• Estrategia:", "¿cómo lo hago inteligentemente?"),
        ("• Impacto:", "¿qué genero?"),
        ("• Dinero:", "¿cómo genero dinero para que mi sustento personal sea sostenible?")
    ]

    col_w = Inches(5.7)
    for idx, (label, desc) in enumerate(items_formula):
        col = 0 if idx < 4 else 1
        row = idx if idx < 4 else idx - 4
        
        x = Inches(0.8) + col * Inches(6.033)
        y = Inches(1.85) + row * Inches(1.25)
        h = Inches(1.05)

        box = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, col_w, h)
        box.fill.solid()
        box.fill.fore_color.rgb = C_SURFACE_DARK
        box.line.color.rgb = RGBColor(60, 70, 90)

        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.15)
        tf.margin_left = Inches(0.3)

        p = tf.paragraphs[0]
        p.text = f"{label} "
        p.font.name = FONT_MAIN
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = C_GOLD

        run = p.add_run()
        run.text = desc
        run.font.name = FONT_MAIN
        run.font.size = Pt(12)
        run.font.bold = False
        run.font.color.rgb = C_TEXT_LIGHT

    # ==========================================
    # SLIDE 9: CONCIENCIA EMPRESARIAL Y MARCA PERSONAL (CLARA)
    # ==========================================
    s9 = prs.slides.add_slide(blank_layout)
    set_bg(s9, C_LIGHT_BG)
    add_header(s9, "Conciencia Empresarial & Marca Personal")

    # Card 1: Texto de toma de decisiones
    box1 = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.75), Inches(11.733), Inches(3.2))
    box1.fill.solid()
    box1.fill.fore_color.rgb = C_CARD_WHITE
    box1.line.color.rgb = C_BORDER_LIGHT

    tf1 = box1.text_frame
    tf1.word_wrap = True
    tf1.margin_top = Inches(0.3)
    tf1.margin_left = tf1.margin_right = Inches(0.4)

    p1 = tf1.paragraphs[0]
    p1.text = "Un emprendedor consciente sabe quién es, tiene una dirección (propósito) y cada día toma decisiones y acciona en función de sus valores y transformando su mentalidad hacia  lo que necesita hacer para llegar un día a la vez hacia crear y sostener un negocio de impacto que le permita vivir con el estilo de vida que desea."
    p1.font.name = FONT_MAIN
    p1.font.size = Pt(12.5)
    p1.font.color.rgb = C_TEXT_DARK

    p2 = tf1.add_paragraph()
    p2.text = "\nLa conciencia empresarial es aprender a tomar mejores decisiones con los elementos que se tienen disponibles pero sin perder de vista los valores los limites la conexion con el cliente y lo que deseo para mi y mi negocio."
    p2.font.name = FONT_MAIN
    p2.font.size = Pt(12.5)
    p2.font.color.rgb = C_TEXT_DARK

    # Card 2: Marca Personal destacada
    box2 = s9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.15), Inches(11.733), Inches(1.6))
    box2.fill.solid()
    box2.fill.fore_color.rgb = C_GOLD_BG
    box2.line.color.rgb = C_GOLD
    box2.line.width = Pt(1.5)

    tf2 = box2.text_frame
    tf2.word_wrap = True
    tf2.margin_top = Inches(0.35)
    tf2.margin_left = tf2.margin_right = Inches(0.4)

    p_mp = tf2.paragraphs[0]
    p_mp.text = "Estas decisiones puestas en acción a diario crean una\nMARCA PERSONAL QUE SE ENTIENDE Y FUNCIONA PARA GENERAR DINERO"
    p_mp.font.name = FONT_MAIN
    p_mp.font.size = Pt(14)
    p_mp.font.bold = True
    p_mp.alignment = PP_ALIGN.CENTER
    p_mp.font.color.rgb = C_GOLD_DARK

    # ==========================================
    # SLIDE 10: MANIFIESTO DE CIERRE (OSCURA DE LUJO)
    # ==========================================
    s10 = prs.slides.add_slide(blank_layout)
    set_bg(s10, C_DARK_NAVY)

    frame10 = s10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.0), Inches(11.333), Inches(5.5))
    frame10.fill.solid()
    frame10.fill.fore_color.rgb = C_SURFACE_DARK
    frame10.line.color.rgb = C_GOLD
    frame10.line.width = Pt(2)

    tf10 = frame10.text_frame
    tf10.word_wrap = True
    tf10.margin_top = Inches(1.0)
    tf10.margin_left = tf10.margin_right = Inches(0.8)

    p_man_tag = tf10.paragraphs[0]
    p_man_tag.text = "EL ENFOQUE DEL CURSO\n"
    p_man_tag.font.name = FONT_MAIN
    p_man_tag.font.size = Pt(12)
    p_man_tag.font.bold = True
    p_man_tag.alignment = PP_ALIGN.CENTER
    p_man_tag.font.color.rgb = C_GOLD

    p_man = tf10.add_paragraph()
    p_man.text = "“Vamos a aprender a construir un negocio sin desconectarnos de nosotras mismas, pero tampoco vamos a romantizarlo porque emprender es crear un negocio que genere renta que te sustente a vos y tu estilo de vida . Vamos a hablar de propósito, dinero, ventas, de mentalidad y autoliderazgo, pero también de estrategias claves y acción consciente para generar negocios con valor de marca e impacto real que atrae dinero.”"
    p_man.font.name = FONT_MAIN
    p_man.font.size = Pt(16.5)
    p_man.font.bold = True
    p_man.alignment = PP_ALIGN.CENTER
    p_man.font.color.rgb = C_TEXT_LIGHT

    # Guardar presentación
    output_dir = r"c:\Users\ecasa\Documents\MODO LIDER\METODOMODOLIDER\clase1_emprender_conscientemente"
    os.makedirs(output_dir, exist_ok=True)
    file_path = os.path.join(output_dir, "Clase_1_Emprender_Conscientemente_MODOLIDER.pptx")
    prs.save(file_path)
    print(f"Presentation saved successfully with exact verbatim text at: {file_path}")

if __name__ == "__main__":
    create_presentation()
